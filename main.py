import openpyxl
from pptx import Presentation
from itertools import groupby

# Load Excel workbook and select worksheet
wb = openpyxl.load_workbook('MOCK-DATA.xlsx')
ws = wb['MOCK-DATA']

# Retrieve data from worksheet
data = []
for row in ws.iter_rows(min_row=2, values_only=True):
    data.append(list(row))

# Sort the data by the first column (car name)
data.sort(key=lambda x: x[0])

# Group the data by the first column (car name)
grouped_data = []
for key, group in groupby(data, lambda x: x[0]):
    grouped_data.append(list(group))

# Count the number of occurrences of the contents of the 1st column
# with respect to the contents of the 2nd column and add the count in the 3rd column
count_dict = {}
for group in grouped_data:
    for row in group:
        key = (row[0], row[1])
        count = count_dict.get(key, 0)
        count_dict[key] = count + 1

# Remove duplicates and update data list with count
updated_data = []
for row in data:
    key = (row[0], row[1])
    count = count_dict.get(key, 0)
    if count > 0:
        row.append(count)
        updated_data.append(row)

# Remove duplicates from updated data list
unique_data = []
for row in updated_data:
    if row not in unique_data:
        unique_data.append(row)

# Open PowerPoint presentation and select slide
prs = Presentation('MOCK-PRESENTATION.pptx')
slide_index = 0
slide = prs.slides[slide_index]

# Select table on slide and update its data
table_name = 'Mal_code'
table = ''
table_found = False
for shape in slide.shapes:
    if shape.name == table_name:
        table = shape.table
        table_found = True
        break

if not table_found:
    print(f"Table '{table_name}' not found on slide.")

else:
    # Update table data using the data from the unique data list
    row_index = 1
    for row in unique_data:
        for j, value in enumerate(row):
            table.cell(row_index, j).text = str(value)
        row_index += 1
    # Remove rows with empty 3rd column
    for i in reversed(range(1, len(table.rows))):
        if not table.cell(i, 2).text.strip():
            tr = table._tbl.tr_lst[i]
            table._tbl.remove(tr)

# Save the updated PowerPoint presentation
prs.save('MOCK-PRESENTATION.pptx')

